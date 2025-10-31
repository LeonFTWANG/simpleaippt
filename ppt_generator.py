import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from config import Config


class PPTGenerator:
    
    def __init__(self, template_path=None):
        """
        初始化 PPT 生成器
        
        Args:
            template_path: 模板文件路径（可选）
        """
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
            self.using_template = True
        else:
            self.prs = Presentation()
            self.using_template = False
            
    def parse_markdown(self, markdown_text):
        """
        解析 Markdown 文本为结构化数据
        
        优化后的结构：
        - ## 标题 → 章节分隔页
        - ### 副标题 → 独立内容页（每个副标题一张幻灯片）
        
        Args:
            markdown_text: Markdown 文本
            
        Returns:
            list: 幻灯片数据列表
        """
        slides_data = []
        current_chapter = None
        current_section = None
        lines = markdown_text.strip().split('\n')
        
        for line in lines:
            line = line.rstrip()
            
            if not line or line.startswith('---'):
                continue
            
            if line.startswith('## '):
                if current_section:
                    slides_data.append(current_section)
                    current_section = None
                
                current_chapter = line[3:].strip()
                slides_data.append({
                    'type': 'chapter',
                    'title': current_chapter,
                    'subtitles': [],
                    'bullets': []
                })
            
            elif line.startswith('### '):
                if current_section:
                    slides_data.append(current_section)
                
                current_section = {
                    'type': 'content',
                    'title': line[4:].strip(),
                    'chapter': current_chapter, 
                    'subtitles': [],
                    'bullets': []
                }
            
            elif line.startswith('- '):
                if current_section:
                    bullet_text = line[2:].strip()
                    current_section['bullets'].append(bullet_text)
        
        if current_section:
            slides_data.append(current_section)
        
        return slides_data
    
    def _get_slide_layout(self, layout_type='content'):
        """
        获取幻灯片布局
        
        Args:
            layout_type: 布局类型 ('title', 'content')
            
        Returns:
            SlideLayout: 幻灯片布局
        """
        layouts = self.prs.slide_layouts
        
        if layout_type == 'title':
            return layouts[0]
        else:
            for i, layout in enumerate(layouts):
                if i == 1 or i == 5: 
                    return layout
            return layouts[1] if len(layouts) > 1 else layouts[0]
    
    def _debug_print_placeholders(self, slide):
        """
        调试方法：打印幻灯片的所有占位符信息
        
        Args:
            slide: 幻灯片对象
        """
        print(f"\n=== 幻灯片占位符调试信息 ===")
        print(f"总共有 {len(slide.placeholders)} 个占位符")
        for i, shape in enumerate(slide.placeholders):
            print(f"占位符 {i}:")
            print(f"  - 类型: {shape.placeholder_format.type}")
            print(f"  - 名称: {shape.name}")
            print(f"  - 是否有 text_frame: {hasattr(shape, 'text_frame')}")
        print("=" * 40)
    
    def add_title_slide(self, title, subtitle=""):
        """
        添加标题幻灯片
        
        Args:
            title: 主标题
            subtitle: 副标题
        """
        slide_layout = self._get_slide_layout('title')
        slide = self.prs.slides.add_slide(slide_layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = title
            self._format_title(slide.shapes.title)
        
        if len(slide.placeholders) > 1 and subtitle:
            slide.placeholders[1].text = subtitle
    
    def add_chapter_slide(self, chapter_title):
        """
        添加章节分隔页
        
        Args:
            chapter_title: 章节标题
        """
        slide_layout = self._get_slide_layout('title')
        slide = self.prs.slides.add_slide(slide_layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = chapter_title
            self._format_title(slide.shapes.title)
    
    def add_content_slide(self, slide_data):
        """
        添加内容幻灯片
        
        Args:
            slide_data: 幻灯片数据字典
        """
        slide_layout = self._get_slide_layout('content')
        slide = self.prs.slides.add_slide(slide_layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = slide_data['title']
            self._format_title(slide.shapes.title)
        
        self._add_content(slide, slide_data)
    
    def _add_content(self, slide, slide_data):
        """
        添加幻灯片内容
        
        优化：更智能地查找和使用占位符，避免创建多余的文本框
        
        Args:
            slide: 幻灯片对象
            slide_data: 内容数据
        """
        content_placeholder = None
        
        for shape in slide.placeholders:
            if hasattr(shape, 'text_frame') and shape.placeholder_format.type == 2:
                content_placeholder = shape
                break
        
        if not content_placeholder:
            for shape in slide.placeholders:
                if hasattr(shape, 'text_frame') and shape.placeholder_format.type == 7:
                    content_placeholder = shape
                    break
        
        if not content_placeholder:
            for shape in slide.placeholders:
                if hasattr(shape, 'text_frame') and shape.placeholder_format.type != 1:  # 排除标题
                    content_placeholder = shape
                    break
        
        if not content_placeholder:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.name.startswith('TextBox'):
                    content_placeholder = shape
                    break
        
        if content_placeholder and hasattr(content_placeholder, 'text_frame'):
            tf = content_placeholder.text_frame
            for paragraph in tf.paragraphs:
                paragraph.clear()
            if not tf.paragraphs:
                tf.text = ''
        else:
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(5)
            shape = slide.shapes.add_textbox(left, top, width, height)
            tf = shape.text_frame
        
        tf.word_wrap = True
        
        # 对于内容页，只添加要点（不再显示副标题，因为副标题已经是幻灯片标题了）
        bullets = slide_data.get('bullets', [])
        
        if bullets:
            p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
            p.text = bullets[0]
            p.level = 0
            p.font.size = Pt(20)
            if not self.using_template:
                p.font.color.rgb = RGBColor(89, 89, 89)
            
            for bullet in bullets[1:]:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = Pt(20)
                if not self.using_template:
                    p.font.color.rgb = RGBColor(89, 89, 89)
    
    def _format_title(self, title_shape):
        """
        格式化标题
        
        Args:
            title_shape: 标题形状对象
        """
        if not self.using_template:
            title_shape.text_frame.paragraphs[0].font.size = Pt(32)
            title_shape.text_frame.paragraphs[0].font.bold = True
            title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(47, 84, 150)
    
    def generate_from_markdown(self, markdown_text, ppt_title="AI 生成的 PPT", ppt_subtitle=""):
        """
        从 Markdown 生成完整的 PPT
        
        优化后的结构：
        1. 总标题页
        2. 章节分隔页（## 标题）
        3. 内容页（### 副标题 + 要点）
        
        Args:
            markdown_text: Markdown 文本
            ppt_title: PPT 总标题
            ppt_subtitle: PPT 副标题
        """
        self.add_title_slide(ppt_title, ppt_subtitle)
        
        slides_data = self.parse_markdown(markdown_text)
        
        for slide_data in slides_data:
            if slide_data.get('type') == 'chapter':
                self.add_chapter_slide(slide_data['title'])
            elif slide_data.get('type') == 'content':
                self.add_content_slide(slide_data)
    
    def save(self, output_path):
        """
        保存 PPT 文件
        
        Args:
            output_path: 输出文件路径
        """
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        self.prs.save(output_path)
        return output_path


def generate_ppt(markdown_content, output_path, template_path=None, title="AI 生成的 PPT", subtitle=""):
    """
    生成 PPT 的便捷函数
    
    Args:
        markdown_content: Markdown 内容
        output_path: 输出路径
        template_path: 模板路径（可选）
        title: PPT 标题
        subtitle: PPT 副标题
        
    Returns:
        str: 生成的文件路径
    """
    generator = PPTGenerator(template_path)
    generator.generate_from_markdown(markdown_content, title, subtitle)
    return generator.save(output_path)

